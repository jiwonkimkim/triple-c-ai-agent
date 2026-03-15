"""
NHN Service QA Portfolio PPT Generator
기존 포트폴리오 디자인(#F0F4FF 배경, #3B7DE8 파란, 화이트 카드)을 재현
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
BG        = RGBColor(0xF0, 0xF4, 0xFF)   # 연라벤더 배경
BLUE      = RGBColor(0x3B, 0x7D, 0xE8)   # 강조 파란
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x8A)   # 진한 파란
NAVY      = RGBColor(0x1E, 0x29, 0x3B)   # 본문 거의 검정
GRAY      = RGBColor(0x64, 0x74, 0x8B)   # 서브 회색
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE= RGBColor(0xDB, 0xEA, 0xFE)   # 뱃지 배경
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
YELLOW    = RGBColor(0xF5, 0x9E, 0x0B)
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HDR = RGBColor(0x3B, 0x7D, 0xE8)
DECO_CIRC = RGBColor(0xC7, 0xD7, 0xFF)   # 장식 원 (연한)

# ── 슬라이드 크기 (16:9 와이드) ─────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── 헬퍼 함수들 ─────────────────────────────────────────────

def add_bg(slide):
    """배경 색상"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_deco_circles(slide):
    """우하단 장식 원"""
    for (l, t, s) in [
        (Inches(11.5), Inches(5.8), Inches(2.5)),
        (Inches(12.2), Inches(6.2), Inches(1.8)),
    ]:
        sh = slide.shapes.add_shape(1, l, t, s, s)
        sh.fill.solid()
        sh.fill.fore_color.rgb = DECO_CIRC
        sh.line.fill.background()

def rect(slide, l, t, w, h, fill=WHITE, line_color=None, line_width=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_width or Pt(1)
    else:
        sh.line.fill.background()
    return sh

def textbox(slide, l, t, w, h, text, size=10, bold=False, color=NAVY,
            align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_section_label(slide, label):
    """좌상단 — 01. SUMMARY 스타일"""
    # 파란 짧은 선
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(0.38), Inches(0.25), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    textbox(slide, Inches(0.82), Inches(0.28), Inches(4), Inches(0.35),
            label, size=9, bold=True, color=BLUE)

def add_headline(slide, text, top=Inches(0.7), size=28):
    textbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.85),
            text, size=size, bold=True, color=NAVY)

def add_subline(slide, text, top=Inches(1.25)):
    textbox(slide, Inches(0.5), top, Inches(12), Inches(0.4),
            text, size=10, color=GRAY)

def card(slide, l, t, w, h, title, body, title_color=NAVY,
         bg=WHITE, border=None, title_size=11, body_size=9):
    rect(slide, l, t, w, h, fill=bg,
         line_color=border or RGBColor(0xE2, 0xE8, 0xF0), line_width=Pt(0.5))
    # 타이틀
    textbox(slide, l+Inches(0.18), t+Inches(0.12), w-Inches(0.3), Inches(0.35),
            title, size=title_size, bold=True, color=title_color)
    # 본문
    textbox(slide, l+Inches(0.18), t+Inches(0.45), w-Inches(0.3), h-Inches(0.55),
            body, size=body_size, color=NAVY)

def badge(slide, l, t, text, bg=LIGHT_BLUE, color=BLUE):
    w = Inches(len(text)*0.11 + 0.3)
    rect(slide, l, t, w, Inches(0.22), fill=bg)
    textbox(slide, l+Inches(0.1), t+Inches(0.02), w, Inches(0.2),
            text, size=8, bold=True, color=color)

def blue_table_row(slide, cols, l, t, w, h, is_header=False):
    col_w = w / len(cols)
    for i, txt in enumerate(cols):
        fill = TABLE_HDR if is_header else (WHITE if i % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFF))
        rect(slide, l + col_w*i, t, col_w, h,
             fill=fill,
             line_color=RGBColor(0xE2, 0xE8, 0xF0))
        textbox(slide, l + col_w*i + Inches(0.1), t + Inches(0.05),
                col_w - Inches(0.15), h - Inches(0.08),
                txt, size=8 if not is_header else 8,
                bold=is_header,
                color=WHITE if is_header else NAVY)

# ════════════════════════════════════════════════════════════
# Page 1 — COVER
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)

# 좌측 레이블
textbox(sl, Inches(0.5), Inches(0.55), Inches(6), Inches(0.3),
        "NHN SERVICE SUPPORT   |   QA ENGINEER PORTFOLIO",
        size=9, bold=True, color=BLUE)

# 메인 헤드라인
for i, line in enumerate(["결함을 설계로 막는", "QA 엔지니어", "김지원입니다"]):
    c = BLUE if i == 1 else NAVY
    sz = 38 if i == 1 else 32
    textbox(sl, Inches(0.5), Inches(1.0)+Inches(i*0.85), Inches(5.8), Inches(0.9),
            line, size=sz, bold=True, color=c)

# 키워드 뱃지
badges = ["멀티 브라우저 E2E 자동화", "API 계약 검증",
          "부하 테스트 설계", "결함 회고 분석", "테스트 자산화"]
bx = Inches(0.5)
for i, b in enumerate(badges):
    bw = Inches(len(b)*0.125 + 0.4)
    rect(sl, bx, Inches(3.85) + (Inches(0.35) if i >= 3 else 0),
         bw, Inches(0.28),
         fill=WHITE,
         line_color=BLUE, line_width=Pt(1))
    textbox(sl, bx+Inches(0.1), Inches(3.87) + (Inches(0.35) if i >= 3 else 0),
            bw, Inches(0.25), b, size=9, color=BLUE)
    bx += bw + Inches(0.12)
    if i == 1:
        bx = Inches(0.5)

# 구분선
ln = rect(sl, Inches(0.5), Inches(4.7), Inches(0.4), Pt(3),
          fill=BLUE)

# 연락처
textbox(sl, Inches(0.5), Inches(4.82), Inches(5.5), Inches(0.28),
        "✉  kjwo516@naver.com", size=10, color=GRAY)
textbox(sl, Inches(0.5), Inches(5.15), Inches(6), Inches(0.28),
        "⚙  Playwright · Postman · nGrinder · Python · GitHub Actions",
        size=9, color=GRAY)

# 우측 카드 1 — Test Experience
rect(sl, Inches(6.0), Inches(0.5), Inches(6.8), Inches(2.9),
     fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
textbox(sl, Inches(6.2), Inches(0.7), Inches(6.3), Inches(0.35),
        "Test Experience", size=13, bold=True, color=NAVY)
for i, (title, body) in enumerate([
    ("RAG 파이프라인 (pytest)",
     "7개 핵심 모듈 단위 검증 및 에러 시나리오 설계"),
    ("Triple C AI 에이전트 (Playwright)",
     "멀티 브라우저 E2E 자동화 + 재발 방지 테스트 자산화"),
]):
    ty = Inches(1.2) + Inches(i*1.1)
    # 체크 아이콘 (파란 원)
    rect(sl, Inches(6.2), ty+Inches(0.05), Inches(0.22), Inches(0.22),
         fill=BLUE)
    textbox(sl, Inches(6.2), ty+Inches(0.05), Inches(0.22), Inches(0.22),
            "✓", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(6.5), ty, Inches(6.1), Inches(0.28),
            title, size=11, bold=True, color=NAVY)
    textbox(sl, Inches(6.5), ty+Inches(0.32), Inches(6.1), Inches(0.35),
            body, size=9, color=GRAY)

# 우측 카드 2 — Automation Practice (파란 배경)
rect(sl, Inches(6.0), Inches(3.55), Inches(6.8), Inches(2.1),
     fill=DARK_BLUE)
textbox(sl, Inches(6.2), Inches(3.7), Inches(6.3), Inches(0.35),
        "Automation Practice", size=12, bold=True, color=WHITE)
for i, line in enumerate([
    "🤖  멀티 브라우저 자동화 — Chromium / Firefox / WebKit",
    "📮  91개 엔드포인트 API 계약 테스트 — Postman Collection",
    "📊  부하 테스트 환경 설계 — nGrinder Docker Compose",
]):
    textbox(sl, Inches(6.2), Inches(4.15)+Inches(i*0.42), Inches(6.3), Inches(0.35),
            line, size=10, color=WHITE)

# ════════════════════════════════════════════════════════════
# Page 2 — 01. SUMMARY
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "01. SUMMARY")
add_headline(sl, "결함을 코드 이전에 설계로 차단하는 QA 엔지니어")

# 좌측 프로필 카드
rect(sl, Inches(0.4), Inches(1.55), Inches(3.5), Inches(5.5),
     fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
textbox(sl, Inches(0.6), Inches(2.3), Inches(3.1), Inches(0.45),
        "김지원", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
textbox(sl, Inches(0.6), Inches(2.78), Inches(3.1), Inches(0.3),
        "NHN Service QA 엔지니어 지원자", size=10, color=BLUE, align=PP_ALIGN.CENTER)
textbox(sl, Inches(0.55), Inches(3.2), Inches(3.2), Inches(2.5),
        "fix: 커밋 20건을 역분석하여\n결함 예방 테스트를 설계하고\n재발 방지 자산으로 전환하는\nQA 엔지니어입니다.",
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 우측 강점 카드 5개
cards_data = [
    ("1. 테스트 자동화 설계",
     "Playwright로 멀티 브라우저 E2E 자동화를 설계하고\n재발 방지 테스트를 GitHub에 자산화했습니다."),
    ("2. 체계적 테스트 설계",
     "ISTQB 기반 경계값 분석·동등 분할·에러 추정 기법을\n실제 20건의 결함 케이스에 적용해 검증했습니다."),
    ("3. 결함 원인 추적",
     "소스 코드 레벨에서 필터 조건·상태 분기·Race Condition을\n추적하여 근본 원인을 도출하고 수정 방향을 제시했습니다."),
    ("4. 성능 취약점 분석",
     "nGrinder Docker 환경을 직접 구성하고\n/api/generate/detail-page의 병목 구간(DB 직렬 쿼리,\nRace Condition)을 코드 레벨에서 식별·문서화했습니다."),
]
for i, (title, body) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    cl = Inches(4.1) + col * Inches(4.4)
    ct = Inches(1.55) + row * Inches(2.0)
    rect(sl, cl, ct, Inches(4.2), Inches(1.8),
         fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
    textbox(sl, cl+Inches(0.2), ct+Inches(0.15), Inches(3.8), Inches(0.35),
            title, size=11, bold=True, color=NAVY)
    textbox(sl, cl+Inches(0.2), ct+Inches(0.5), Inches(3.8), Inches(1.1),
            body, size=9, color=NAVY)

# 5번 wide 카드
rect(sl, Inches(4.1), Inches(5.6), Inches(8.8), Inches(1.35),
     fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
textbox(sl, Inches(4.3), Inches(5.72), Inches(8.4), Inches(0.35),
        "5. 테스트 자산 문서화", size=11, bold=True, color=NAVY)
textbox(sl, Inches(4.3), Inches(6.08), Inches(8.4), Inches(0.7),
        "결함 회고 분석(RETROSPECTIVE_TEST_ANALYSIS.md), 성능 테스트 보고서,\nAPI 계약 테스트 컬렉션(91개)을 GitHub에 자산화하여 재현 가능한 QA 산출물로 관리합니다.",
        size=9, color=NAVY)

# ════════════════════════════════════════════════════════════
# Page 3 — 02. SKILL
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "02. SKILL")
add_headline(sl, "검증에 직접 활용한 테스트 도구 스택")
add_subline(sl, "각 도구의 도입 목적과 실제 검증 성과 중심으로 기술합니다.")

skill_cards = [
    ("🎭  Playwright", "멀티 브라우저 E2E 자동화",
     "Microsoft 개발. Chromium·Firefox·WebKit 3개 브라우저를 단일 코드로 동시 검증.\nDB 없이 JWT를 직접 발급(compat-auth.setup.ts)하여 독립 실행 가능한 테스트 환경을 설계했습니다.",
     "✔ 이미지 전용 메시지 누락 결함 — TC-1.1~1.4 재발 방지 자동화\n✔ fontSize 경계값(12·48·65·80) BVA 테스트 설계 및 실행\n✔ SSE 스트리밍·API 에러 5개 시나리오 크로스브라우저 통과"),
    ("📮  Postman", "API 계약 테스트 · 91개 엔드포인트",
     "GUI 기반 REST API 검증 도구.\n요청/응답 스키마 계약을 Collection으로 관리하고\n팀 전체가 공유·재사용 가능한 API 테스트 자산을 구축합니다.",
     "✔ 14개 폴더, 91개 요청 컬렉션 설계 및 자산화\n✔ 401 Unauthorized · 400 Bad Request 음수 시나리오 검증\n✔ SSE 스트리밍 · multipart/form-data 엣지 케이스 포함"),
    ("📊  nGrinder", "NHN 오픈소스 부하 테스트 설계",
     "NHN(구 네이버)이 개발한 오픈소스 성능 테스트 플랫폼.\nGroovy 스크립트로 실제 유저 시나리오를 구현하고\nTPS · Peak TPS · MTT · 에러율을 측정합니다.",
     "✔ Docker Compose로 Controller + Agent 로컬 환경 구성\n✔ 4-Step 시나리오 (CSRF→로그인→프로젝트→AI 생성) 설계\n✔ SLA 10s 기준 응답시간 검증 로직 및 성능 보고서 자산화\n✔ AI 호출·DB 직렬 쿼리·Race Condition 병목 구간 분석"),
]

for i, (title, tag, value, result) in enumerate(skill_cards):
    cl = Inches(0.4) + i * Inches(4.28)
    ct = Inches(1.55)
    rect(sl, cl, ct, Inches(4.08), Inches(4.05),
         fill=WHITE, line_color=BLUE, line_width=Pt(1.5))
    # 도구 타이틀
    textbox(sl, cl+Inches(0.2), ct+Inches(0.15), Inches(3.7), Inches(0.38),
            title, size=13, bold=True, color=NAVY)
    # 태그 뱃지
    rect(sl, cl+Inches(0.2), ct+Inches(0.55), Inches(3.5), Inches(0.22),
         fill=LIGHT_BLUE)
    textbox(sl, cl+Inches(0.25), ct+Inches(0.56), Inches(3.4), Inches(0.2),
            tag, size=8, bold=True, color=BLUE)
    # 구분선
    rect(sl, cl+Inches(0.2), ct+Inches(0.85), Inches(3.5), Pt(0.5), fill=RGBColor(0xE2,0xE8,0xF0))
    # 활용 가치
    textbox(sl, cl+Inches(0.2), ct+Inches(0.9), Inches(3.7), Inches(0.22),
            "활용 가치", size=8, bold=True, color=GRAY)
    textbox(sl, cl+Inches(0.2), ct+Inches(1.1), Inches(3.7), Inches(1.0),
            value, size=8.5, color=NAVY)
    # 구분선
    rect(sl, cl+Inches(0.2), ct+Inches(2.15), Inches(3.5), Pt(0.5), fill=RGBColor(0xE2,0xE8,0xF0))
    # 검증 성과
    textbox(sl, cl+Inches(0.2), ct+Inches(2.2), Inches(3.7), Inches(0.22),
            "검증 성과", size=8, bold=True, color=GRAY)
    textbox(sl, cl+Inches(0.2), ct+Inches(2.4), Inches(3.7), Inches(1.45),
            result, size=8.5, color=NAVY)

# 하단 보조 스킬 표
table_top = Inches(5.72)
headers = ["도구 / 언어", "활용 내용"]
rows = [
    ["Python · pytest", "RAG 파이프라인 7개 모듈 단위 테스트, JUnit XML 보고서"],
    ["GitHub Actions", "Push·PR 트리거 CI 파이프라인 + Quality Gate 설정"],
    ["Docker Compose", "개발-프로덕션 동일 환경 구성, nGrinder 인프라 설계"],
    ["TypeScript · Next.js", "프로덕션 코드 기반 테스트 환경 구성 경험"],
]
col_widths = [Inches(2.5), Inches(9.8)]
tl = Inches(0.4)
blue_table_row(sl, headers, tl, table_top, Inches(12.53), Inches(0.32), is_header=True)
for ri, row in enumerate(rows):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFF)
    for ci, cell in enumerate(row):
        cx = tl + (col_widths[0] if ci == 1 else Inches(0))
        rect(sl, cx, table_top + Inches(0.32) + ri * Inches(0.3),
             col_widths[ci], Inches(0.3),
             fill=bg, line_color=RGBColor(0xE2, 0xE8, 0xF0))
        textbox(sl, cx+Inches(0.1), table_top + Inches(0.34) + ri*Inches(0.3),
                col_widths[ci]-Inches(0.15), Inches(0.28),
                cell, size=8.5, bold=(ci==0), color=NAVY)

# ════════════════════════════════════════════════════════════
# Page 4 — 03. TROUBLESHOOTING
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "03. TROUBLESHOOTING")
add_headline(sl, "핵심 사례: 에러 추정 기법으로 탐지한 UI 상태 누락 결함")

# 좌측 — 결함 흐름 다이어그램
rect(sl, Inches(0.4), Inches(1.55), Inches(5.2), Inches(5.5),
     fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
textbox(sl, Inches(0.6), Inches(1.65), Inches(4.8), Inches(0.3),
        "Defect Flow Analysis", size=10, bold=True, color=NAVY)

# 흐름 박스 3개
for i, (icon, title, body, fill_c, txt_c) in enumerate([
    ("📱", "Chat UI (Next.js)",
     "message.content = \"\"\nmessage.attachments = [url]",
     WHITE, NAVY),
    ("❌", "use-chat.ts 메시지 필터 (결함)",
     "prev.filter(m => m.content.trim())\n↑ content=\"\"이면 attachments\n  유무 무관하게 메시지 제거",
     RGBColor(0xFF,0xEE,0xEE), RGBColor(0xDC,0x26,0x26)),
    ("📭", "대화창에서 메시지 소실",
     "사용자는 전송 실패로 오인",
     WHITE, NAVY),
]):
    by = Inches(2.1) + i * Inches(1.35)
    rect(sl, Inches(0.6), by, Inches(4.8), Inches(1.1),
         fill=fill_c, line_color=RGBColor(0xCC,0xCC,0xCC))
    textbox(sl, Inches(0.8), by+Inches(0.08), Inches(4.4), Inches(0.3),
            f"{icon}  {title}", size=10, bold=True, color=txt_c)
    textbox(sl, Inches(0.8), by+Inches(0.4), Inches(4.4), Inches(0.6),
            body, size=9, color=txt_c)
    if i < 2:
        textbox(sl, Inches(2.7), by+Inches(1.12), Inches(0.5), Inches(0.25),
                "▼", size=12, color=BLUE, align=PP_ALIGN.CENTER)

# 코드 카드 (다크)
rect(sl, Inches(0.6), Inches(6.1), Inches(4.8), Inches(0.8),
     fill=RGBColor(0x1E, 0x29, 0x3B))
textbox(sl, Inches(0.75), Inches(6.12), Inches(4.5), Inches(0.35),
        "❌  prev.filter(m => m.content.trim())",
        size=8, color=RGBColor(0xFF,0x6B,0x6B))
textbox(sl, Inches(0.75), Inches(6.47), Inches(4.5), Inches(0.35),
        "✅  return hasContent || hasAttachments",
        size=8, color=RGBColor(0x6B,0xFF,0x9E))

# 우측 — 4단계 워크플로우
steps = [
    ("🔍", "1. 현상 (Bug Detection)", BLUE,
     "이미지만 첨부한 메시지 전송 후 AI 응답 수신 시\n해당 메시지가 대화창에서 즉시 소실됨.\n사용자는 전송 성공 여부를 확인할 수 없는 상태."),
    ("🧠", "2. 원인 분석 — 에러 추정(Error Guessing) 기법 적용", YELLOW,
     "\"content가 빈 문자열인 경우 필터에서 제거될 수 있다\"는\n에러 추정 가설로 접근 → use-chat.ts의 메시지 필터 조건 확인.\ncontent.trim()만으로 필터링하는 로직이 attachments를 가진\n메시지를 동일하게 제거하는 것을 확인. → 조건 누락이 근본 원인."),
    ("🔧", "3. 해결 (Resolution)", GREEN,
     "hasContent || hasAttachments 복합 조건으로 수정.\ncontent가 없더라도 attachments가 존재하면 메시지를 유지.\n수정 커밋: 12146b1"),
    ("🛡", "4. 예방 — Playwright 자동화로 재발 차단", BLUE,
     "TC-1.1: 이미지 전용 첨부 → AI 응답 후 메시지 visible 검증\nTC-1.2: DB 재조회 경로(content=\"\" + attachments=[url]) 검증\nTC-1.3: 텍스트+이미지 혼합 메시지 회귀 검증\nTC-1.4: 빈 메시지 정상 필터 검증\n→ tests/e2e/bug-prevention.spec.ts 자산화 완료"),
]
for i, (icon, title, accent, body) in enumerate(steps):
    sy = Inches(1.55) + i * Inches(1.48)
    rect(sl, Inches(5.8), sy, Inches(7.1), Inches(1.38),
         fill=WHITE, line_color=RGBColor(0xE2, 0xE8, 0xF0))
    # 아이콘 원
    rect(sl, Inches(5.95), sy+Inches(0.1), Inches(0.4), Inches(0.4),
         fill=RGBColor(0xDB,0xEA,0xFE))
    textbox(sl, Inches(5.95), sy+Inches(0.1), Inches(0.4), Inches(0.4),
            icon, size=12, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(6.42), sy+Inches(0.1), Inches(6.3), Inches(0.32),
            title, size=10, bold=True, color=NAVY)
    textbox(sl, Inches(6.42), sy+Inches(0.45), Inches(6.3), Inches(0.85),
            body, size=8.5, color=NAVY)

# ════════════════════════════════════════════════════════════
# Page 5 — 04. QA PROCESS (신규)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "04. QA PROCESS")
add_headline(sl, "20건의 결함 회고 분석 — ISTQB 기반 설계 기법 실전 적용")
add_subline(sl, "fix: 커밋 20건을 QA 엔지니어 관점에서 역분석하여 \"코딩 전에 잡을 수 있었던 테스트\"를 설계하고 GitHub에 자산화했습니다.")

# 좌측 — ISTQB 기법 표
tl2 = Inches(0.4)
tt2 = Inches(1.62)
headers2 = ["ISTQB 기법", "정의", "적용 사례 (Triple C)"]
widths2 = [Inches(2.2), Inches(2.0), Inches(4.8)]
# 헤더
x = tl2
for ci, (h, w) in enumerate(zip(headers2, widths2)):
    rect(sl, x, tt2, w, Inches(0.35), fill=TABLE_HDR,
         line_color=RGBColor(0xE2,0xE8,0xF0))
    textbox(sl, x+Inches(0.1), tt2+Inches(0.07), w-Inches(0.15), Inches(0.28),
            h, size=8, bold=True, color=WHITE)
    x += w

rows2 = [
    ("에러 추정\n(Error Guessing)", "경험 기반 결함 예상",
     "이미지 전용 메시지 소실 — content=\"\" 조건 누락 추정"),
    ("경계값 분석\n(BVA)", "입력 범위 경계 테스트",
     "fontSize {12, 48, 49, 65, 80} 5개 경계값 설계"),
    ("동등 분할\n(EP)", "동등 입력 그룹화",
     "채팅 옵션 선택값 유효/무효 파티션 분리"),
    ("상태 전이\n(State Transition)", "시스템 상태 변화 추적",
     "Lobby→정보수집→기획→생성→피드백 6단계 전이 검증"),
    ("CRUD 완전성", "생성/조회/수정/삭제 망라",
     "프로젝트·브랜드·템플릿 API 91개 계약 테스트"),
]
for ri, row in enumerate(rows2):
    x = tl2
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFF)
    rh = Inches(0.68)
    for ci, (cell, w) in enumerate(zip(row, widths2)):
        rect(sl, x, tt2+Inches(0.35)+ri*rh, w, rh,
             fill=bg, line_color=RGBColor(0xE2,0xE8,0xF0))
        textbox(sl, x+Inches(0.1), tt2+Inches(0.37)+ri*rh,
                w-Inches(0.15), rh-Inches(0.08),
                cell, size=8, bold=(ci==0), color=NAVY)
        x += w

# 우측 — 성과 수치 카드 3개
nums = [("20", "fix: 커밋\n역분석"), ("5", "ISTQB 기법\n실전 적용"), ("9", "Playwright\n테스트 케이스")]
for i, (num, label) in enumerate(nums):
    nx = Inches(9.5) + i * Inches(1.28)
    rect(sl, nx, Inches(1.62), Inches(1.15), Inches(1.5),
         fill=WHITE, line_color=BLUE, line_width=Pt(1.5))
    textbox(sl, nx, Inches(1.72), Inches(1.15), Inches(0.65),
            num, size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    textbox(sl, nx, Inches(2.38), Inches(1.15), Inches(0.65),
            label, size=8, color=GRAY, align=PP_ALIGN.CENTER)

# 우측 아래 — 적용 커밋 분포 미니 카드
rect(sl, Inches(9.5), Inches(3.28), Inches(3.4), Inches(2.5),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
textbox(sl, Inches(9.65), Inches(3.38), Inches(3.1), Inches(0.3),
        "결함 카테고리 분포", size=10, bold=True, color=NAVY)
cats = [
    ("UI/UX 결함", "8건", 0.60),
    ("상태 관리", "5건", 0.38),
    ("API 통신", "4건", 0.30),
    ("비즈니스 로직", "3건", 0.23),
]
for i, (cat, cnt, ratio) in enumerate(cats):
    cy = Inches(3.75) + i * Inches(0.42)
    textbox(sl, Inches(9.65), cy, Inches(1.5), Inches(0.28),
            cat, size=8.5, color=NAVY)
    textbox(sl, Inches(11.4), cy, Inches(0.5), Inches(0.28),
            cnt, size=8.5, bold=True, color=BLUE)
    # 바 배경
    rect(sl, Inches(9.65), cy+Inches(0.28), Inches(3.0), Inches(0.1),
         fill=RGBColor(0xE2,0xE8,0xF0))
    rect(sl, Inches(9.65), cy+Inches(0.28), Inches(3.0*ratio), Inches(0.1),
         fill=BLUE)

# 하단 산출물 카드 (파란 배경)
rect(sl, Inches(0.4), Inches(6.1), Inches(12.53), Inches(1.1),
     fill=DARK_BLUE)
textbox(sl, Inches(0.6), Inches(6.22), Inches(6.0), Inches(0.3),
        "📄  docs/RETROSPECTIVE_TEST_ANALYSIS.md", size=9, bold=True, color=WHITE)
textbox(sl, Inches(0.6), Inches(6.55), Inches(5.8), Inches(0.35),
        "20개 사례별 · Playwright 아이디어 · 설계 기법 매핑 · GitHub 자산화",
        size=8.5, color=RGBColor(0xBF,0xDB,0xFF))
textbox(sl, Inches(6.6), Inches(6.22), Inches(5.9), Inches(0.3),
        "🧪  tests/e2e/bug-prevention.spec.ts", size=9, bold=True, color=WHITE)
textbox(sl, Inches(6.6), Inches(6.55), Inches(5.9), Inches(0.35),
        "Case #1 (이미지 메시지 누락) · Case #14 (fontSize BVA) 실행 가능 구현체",
        size=8.5, color=RGBColor(0xBF,0xDB,0xFF))

# ════════════════════════════════════════════════════════════
# Page 6 — 05. PROJECT (Triple C)
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "05. PROJECT 02")
add_headline(sl, "Triple C AI 에이전트 — 통합/E2E 자동화 및 성능 테스트 설계")

proj_cards = [
    ("1. 프로덕션 동일 코드 E2E 검증",
     "실제 서비스 로직 기반 테스트 환경 구성.\nDB 없이 JWT를 직접 발급하는 compat-auth.setup.ts를 설계하여\n외부 의존 없이 독립 실행 가능한 인증 환경을 구축했습니다."),
    ("2. 결함 회고 기반 재발 방지 자동화",
     "fix: 커밋 20건 역분석 → ISTQB 기법 적용 →\n재발 방지 Playwright 테스트 9개 케이스 설계 및 자산화.\n단순 버그 수정을 넘어 테스트 자산으로 전환했습니다."),
    ("3. CI/CD 품질 자동화 환경",
     "GitHub Actions Push·PR 트리거 CI 파이프라인 구성.\nDocker Compose 통합 환경으로 개발-프로덕션 일관성 확보.\n테스트 통과 시에만 배포되도록 Quality Gate 정책 설계."),
]
for i, (title, body) in enumerate(proj_cards):
    cl = Inches(0.4) + i*Inches(4.28)
    ct = Inches(1.55)
    rect(sl, cl, ct, Inches(4.1), Inches(2.1),
         fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
    textbox(sl, cl+Inches(0.2), ct+Inches(0.15), Inches(3.7), Inches(0.35),
            title, size=11, bold=True, color=NAVY)
    textbox(sl, cl+Inches(0.2), ct+Inches(0.55), Inches(3.7), Inches(1.35),
            body, size=9, color=NAVY)

# CI/CD 플로우
rect(sl, Inches(0.4), Inches(3.82), Inches(12.53), Inches(1.8),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
for i, (icon, label, c) in enumerate([
    ("⬆", "Code Push / PR", BLUE),
    ("→", "", GRAY),
    ("🔨", "Build & Test", BLUE),
    ("→", "", GRAY),
    ("🛡", "Quality Gate", GREEN),
    ("→", "", GRAY),
    ("🚀", "Deploy", BLUE),
]):
    bx = Inches(1.0) + i * Inches(1.65)
    if icon in ("→",):
        textbox(sl, bx, Inches(4.55), Inches(0.5), Inches(0.5),
                icon, size=18, color=GRAY, align=PP_ALIGN.CENTER)
    else:
        rect(sl, bx, Inches(4.0), Inches(1.3), Inches(0.7),
             fill=RGBColor(0xDB,0xEA,0xFE))
        textbox(sl, bx, Inches(4.05), Inches(1.3), Inches(0.65),
                f"{icon}\n{label}", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 산출물 (Deliverables)
rect(sl, Inches(0.4), Inches(5.75), Inches(12.53), Inches(1.35),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
textbox(sl, Inches(0.6), Inches(5.82), Inches(2.0), Inches(0.3),
        "산출물 (Deliverables)", size=10, bold=True, color=NAVY)
deliverables = [
    "✅ Playwright E2E 9개 TC (멀티 브라우저 통과 검증)",
    "✅ Postman Collection 91개 엔드포인트 계약 테스트",
    "✅ nGrinder 부하 테스트 스크립트 + Docker 환경",
    "✅ 성능 테스트 보고서 (병목 분석 + 개선 제안 4건)",
    "✅ 결함 회고 분석 문서 (20건, ISTQB 기법 매핑)",
]
for i, d in enumerate(deliverables):
    col = i % 3
    row = i // 3
    textbox(sl, Inches(0.6)+col*Inches(4.1), Inches(6.18)+row*Inches(0.38),
            Inches(3.9), Inches(0.35), d, size=8.5, color=NAVY)

# ════════════════════════════════════════════════════════════
# Page 7 — 06. FUTURE PLAN
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "06. FUTURE PLAN")
add_headline(sl, "지금 설계하고 있는 다음 검증 과제")
add_subline(sl, "NHN Service 도메인 역량 확보를 위해 진행 중인 3가지 과제")

future_cards = [
    ("🎮", "한게임 버그 리포트 10건", "🟡 IN PROGRESS", YELLOW,
     "과제\n한게임 웹보드 및 모바일 게임을 직접 플레이하며\n정식 버그 리포트 10건을 작성합니다.\n\n목적\nNHN Service 핵심 도메인인 게임 서비스에서\n실제 결함을 탐지하고 리포팅하는 역량을 검증합니다.\n\n리포트 형식\n제목 / 환경(OS·브라우저·디바이스) / 재현 단계 /\n기대 결과 vs 실제 결과 / 심각도 / 스크린샷\n\n산출물 목표\nGitHub: nhn-game-qa-report 저장소 공개"),
    ("📊", "nGrinder 성능 임계치 측정", "🟡 IN PROGRESS", YELLOW,
     "과제\nnGrinder Vuser 1→5→10→20 단계별 부하 테스트 실행 및\nTPS · Peak TPS · MTT · 에러율 실측값을 측정합니다.\n\n목적\n코드 분석으로 식별한 병목 구간(DB 직렬 쿼리,\nRace Condition)이 실제 부하 환경에서\n임계점을 어디서 만드는지 수치로 검증합니다.\n\n산출물\nPERFORMANCE_TEST_REPORT.md 실측값 완성\n개선 제안 Before/After 검증"),
    ("📘", "ISTQB CTFL 취득 준비", "⚪ PLANNED", GRAY,
     "과제\nISTQB Foundation Level 자격 취득을 통해\n테스트 설계 기법의 이론적 기반을 공식 검증합니다.\n\n현재 상태\n이미 경계값 분석·동등 분할·에러 추정·상태 전이 기법을\n실제 프로젝트 20건 결함에 적용 완료.\n교재 기반 이론 보완 후 응시 예정.\n\n기대 효과\n검증된 이론 기반으로 NHN Service 실무에서\n체계적인 테스트 케이스 설계 역량 발휘"),
]

for i, (icon, title, status, sc, body) in enumerate(future_cards):
    cl = Inches(0.4) + i * Inches(4.28)
    ct = Inches(1.6)
    rect(sl, cl, ct, Inches(4.1), Inches(5.2),
         fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
    # 아이콘 원
    rect(sl, cl+Inches(0.2), ct+Inches(0.15), Inches(0.5), Inches(0.5),
         fill=LIGHT_BLUE)
    textbox(sl, cl+Inches(0.2), ct+Inches(0.15), Inches(0.5), Inches(0.5),
            icon, size=16, align=PP_ALIGN.CENTER)
    textbox(sl, cl+Inches(0.8), ct+Inches(0.22), Inches(3.1), Inches(0.35),
            title, size=11, bold=True, color=NAVY)
    # 상태 뱃지
    rect(sl, cl+Inches(0.8), ct+Inches(0.58), Inches(1.5), Inches(0.22),
         fill=LIGHT_BLUE if sc == YELLOW else RGBColor(0xF1,0xF5,0xF9))
    textbox(sl, cl+Inches(0.85), ct+Inches(0.59), Inches(1.4), Inches(0.2),
            status, size=8, bold=True, color=sc)
    # 구분선
    rect(sl, cl+Inches(0.2), ct+Inches(0.9), Inches(3.7), Pt(0.5),
         fill=RGBColor(0xE2,0xE8,0xF0))
    textbox(sl, cl+Inches(0.2), ct+Inches(0.98), Inches(3.7), Inches(4.0),
            body, size=8.5, color=NAVY)

# 하단 인용 카드
rect(sl, Inches(0.4), Inches(6.95), Inches(12.53), Inches(0.4),
     fill=DARK_BLUE)
textbox(sl, Inches(0.6), Inches(6.98), Inches(12.2), Inches(0.32),
        "\"  NHN Service의 게임·커머스·클라우드 서비스 품질을 검증하는 QA 엔지니어로서, 도메인 지식의 공백을 직접 탐지하고 체계적으로 채워나가겠습니다.  \"",
        size=9, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# Page 8 — 07. CLOSING
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_deco_circles(sl)
add_section_label(sl, "07. CLOSING")
add_headline(sl, "QA 엔지니어로서의 3가지 실천 원칙")

# 인용구
rect(sl, Inches(1.5), Inches(1.65), Inches(10.3), Inches(0.75),
     fill=WHITE, line_color=BLUE, line_width=Pt(1.5))
textbox(sl, Inches(1.7), Inches(1.75), Inches(9.9), Inches(0.55),
        "\"결함을 발견하는 것보다 설계 단계에서 예방하는 것이\n더 가치 있는 QA 엔지니어가 되겠습니다.\"",
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 3개 카드
closing_cards = [
    ("🛡", "결함 예방 중심",
     "요구사항과 코드를 함께 분석하여\n테스트 설계 단계에서 결함 가능성을 식별합니다.\nISTQB 기법 기반의 체계적 설계로 초기 결함을 차단합니다."),
    ("🐳", "재현성 우선",
     "환경 표준화(Docker Compose)와\n독립 실행 가능한 테스트 자산(compat-auth, fixture)으로\n언제 어디서나 동일하게 재현 가능한 검증 체계를 구축합니다."),
    ("🔄", "지속적 자산화",
     "수행한 테스트를 문서와 코드로 자산화하여\n팀이 공유하고 재사용 가능한 QA 인프라를 만들어갑니다.\nCI 품질 게이트와 연결해 지속적으로 품질을 측정합니다."),
]
for i, (icon, title, body) in enumerate(closing_cards):
    cl = Inches(0.8) + i * Inches(4.0)
    ct = Inches(2.55)
    rect(sl, cl, ct, Inches(3.7), Inches(3.3),
         fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
    # 아이콘 원형
    rect(sl, cl+Inches(1.3), ct+Inches(0.25), Inches(1.1), Inches(1.1),
         fill=LIGHT_BLUE)
    textbox(sl, cl+Inches(1.3), ct+Inches(0.25), Inches(1.1), Inches(1.1),
            icon, size=28, align=PP_ALIGN.CENTER)
    textbox(sl, cl+Inches(0.2), ct+Inches(1.5), Inches(3.3), Inches(0.38),
            title, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    textbox(sl, cl+Inches(0.2), ct+Inches(1.92), Inches(3.3), Inches(1.15),
            body, size=9, color=NAVY, align=PP_ALIGN.CENTER)

# 하단 연락처 카드
rect(sl, Inches(2.5), Inches(6.15), Inches(8.3), Inches(0.98),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0))
textbox(sl, Inches(3.0), Inches(6.4), Inches(2.5), Inches(0.45),
        "김지원", size=18, bold=True, color=NAVY)
textbox(sl, Inches(5.2), Inches(6.47), Inches(2.0), Inches(0.35),
        "QA Engineer", size=11, color=GRAY)
# 구분선
rect(sl, Inches(7.1), Inches(6.4), Pt(0.5), Inches(0.45),
     fill=RGBColor(0xE2,0xE8,0xF0))
textbox(sl, Inches(7.25), Inches(6.4), Inches(3.0), Inches(0.45),
        "✉  kjwo516@naver.com", size=11, color=NAVY)

# ── 저장 ────────────────────────────────────────────────────
output = "/Users/wonricordo/Downloads/nhn_김지원_qa_portfolio_v2.pptx"
prs.save(output)
print(f"✅ 저장 완료: {output}")
